from pathlib import Path


def parse_pptx(filepath: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx requis: pip install python-pptx")

    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)
